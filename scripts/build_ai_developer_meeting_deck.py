#!/usr/bin/env python3
"""Build the 10-slide AI developer meeting presentation.

Dependencies are intentionally kept outside the project:
  PYTHONPATH=/tmp/vive_ppt_deps python3 scripts/build_ai_developer_meeting_deck.py
"""

from __future__ import annotations

import math
import sys
from pathlib import Path

try:
    import qrcode
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - developer-facing message
    raise SystemExit(
        "Missing build dependencies. Run with PYTHONPATH=/tmp/vive_ppt_deps."
    ) from exc


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
OUT_FILE = OUT_DIR / "ai-developer-meeting-vive-with-gemini-10slides.pptx"
QR_FILE = Path("/tmp/vive-chapter-5-qr.png")

PUBLIC_URL = (
    "https://tako-chan0511.github.io/vive-with-gemini/"
    "ai-agile-vive-with-gemini-5.html"
)

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_JP = "Yu Gothic"
FONT_MONO = "Consolas"

NAVY = "173B5E"
NAVY_DARK = "0F2B46"
BLUE = "2F80ED"
BLUE_LIGHT = "EAF4FF"
ORANGE = "F2994A"
ORANGE_LIGHT = "FFF2E2"
GREEN = "2E8B57"
GREEN_LIGHT = "EAF7F0"
PURPLE = "7B61A8"
PURPLE_LIGHT = "F2ECFA"
RED = "D64545"
RED_LIGHT = "FDECEC"
INK = "1F2937"
MUTED = "536476"
PAPER = "F6F8FB"
WHITE = "FFFFFF"
LINE = "D7E0EA"
GRAY_LIGHT = "EEF2F6"
GOLD = "F5C451"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def set_background(slide, color: str = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str | None = LINE,
    radius: bool = True,
    line_width: float = 1.2,
    rotation: float = 0,
):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.rotation = rotation
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_circle(
    slide,
    x: float,
    y: float,
    d: float,
    *,
    fill: str,
    line: str | None = None,
    line_width: float = 1.2,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    font: str = FONT_JP,
    margin: float = 0.04,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    runs: list[tuple[str, float, str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    for value, size, color, bold in runs:
        run = p.add_run()
        run.text = value
        run.font.name = FONT_JP
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str = "",
    *,
    accent: str = BLUE,
    fill: str = WHITE,
    title_size: float = 18,
    body_size: float = 14,
    number: str | None = None,
    center: bool = False,
):
    add_box(slide, x, y, w, h, fill=fill, line=accent, line_width=1.5)
    add_box(slide, x, y, 0.10, h, fill=accent, line=None, radius=True)
    title_x = x + 0.22
    title_w = w - 0.42
    if number is not None:
        add_circle(slide, x + 0.20, y + 0.18, 0.46, fill=accent)
        add_text(
            slide,
            number,
            x + 0.20,
            y + 0.18,
            0.46,
            0.46,
            size=16,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        title_x = x + 0.78
        title_w = w - 0.98
    align = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    add_text(
        slide,
        title,
        title_x,
        y + 0.12,
        title_w,
        0.55,
        size=title_size,
        color=INK,
        bold=True,
        align=align,
    )
    if body:
        add_text(
            slide,
            body,
            x + 0.28,
            y + 0.70,
            w - 0.50,
            h - 0.82,
            size=body_size,
            color=INK,
            align=align,
            valign=MSO_ANCHOR.TOP,
            line_spacing=1.05,
        )


def add_chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.36,
    *,
    fill: str = GRAY_LIGHT,
    color: str = MUTED,
    line: str | None = None,
    size: float = 12,
    bold: bool = True,
):
    add_box(slide, x, y, w, h, fill=fill, line=line, line_width=0.9)
    add_text(
        slide,
        text,
        x + 0.05,
        y,
        w - 0.10,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
    )


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = NAVY,
    width: float = 2.2,
    dashed: bool = False,
    head: float = 0.18,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    tip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(x2 - head / 2),
        Inches(y2 - head / 2),
        Inches(head),
        Inches(head),
    )
    tip.rotation = angle + 90
    tip.fill.solid()
    tip.fill.fore_color.rgb = rgb(color)
    tip.line.fill.background()
    return connector


def add_chevron(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = BLUE,
    rotation: float = 0,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.rotation = rotation
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    return shape


def add_header(slide, section: str, title: str, page: int) -> None:
    add_box(slide, 0, 0, SLIDE_W, 0.72, fill=NAVY, line=None, radius=False)
    section_size = 10 if len(section) > 6 else 12
    add_chip(
        slide,
        section,
        0.42,
        0.17,
        1.16,
        0.38,
        fill=WHITE,
        color=NAVY,
        size=section_size,
    )
    add_text(
        slide,
        title,
        1.68,
        0.07,
        10.75,
        0.58,
        size=25,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        f"{page:02d}",
        12.48,
        0.10,
        0.43,
        0.48,
        size=13,
        color="D9E8F5",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide, section: str, message: str) -> None:
    add_box(slide, 0.40, 6.87, 12.53, 0.45, fill=NAVY, line=None)
    add_chip(
        slide,
        section,
        0.58,
        6.93,
        1.08,
        0.33,
        fill=ORANGE,
        color=WHITE,
        size=11,
    )
    add_text(
        slide,
        message,
        1.88,
        6.88,
        10.75,
        0.43,
        size=15,
        color=WHITE,
        bold=True,
    )


def add_note(slide, note: str) -> None:
    slide.notes_slide.notes_text_frame.text = note.strip()


def slide_1(prs: Presentation, qr_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY_DARK)

    add_chip(
        slide,
        "AI開発事例共有会｜2026.09.11",
        0.72,
        0.55,
        3.65,
        0.42,
        fill=BLUE,
        color=WHITE,
        size=13,
    )
    add_text(
        slide,
        "AIは「道具」でなく\n「相棒」",
        0.72,
        1.25,
        6.15,
        1.78,
        size=35,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        "― アジャイルなAI駆動開発の実践事例 ―",
        0.75,
        3.05,
        5.85,
        0.55,
        size=19,
        color="D9E8F5",
        bold=True,
    )
    add_box(slide, 0.75, 3.86, 5.52, 0.92, fill="234B6D", line="47769A")
    add_text(
        slide,
        "AIが提案し、人が決める。\n短いフィードバックループを一緒に回す。",
        1.02,
        3.96,
        4.98,
        0.70,
        size=18,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_chip(slide, "64歳から学び直し", 0.75, 5.20, 1.80, fill=WHITE, color=NAVY)
    add_chip(slide, "AI対話で26アプリ", 2.68, 5.20, 1.80, fill=WHITE, color=NAVY)
    add_chip(slide, "現在：閉域MaaS", 4.61, 5.20, 1.66, fill=WHITE, color=NAVY)
    add_text(
        slide,
        "原 桂介｜ビートテック株式会社 九州支店",
        0.76,
        6.05,
        5.75,
        0.55,
        size=17,
        color=WHITE,
        bold=True,
    )

    # Native-shape key visual: Human and AI connected by a dialogue loop.
    add_circle(slide, 8.13, 1.10, 3.85, fill="204B70", line="5B82A2", line_width=1.5)
    add_circle(slide, 8.67, 2.05, 1.18, fill=ORANGE)
    add_text(slide, "人", 8.67, 2.05, 1.18, 1.18, size=25, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_circle(slide, 10.27, 2.05, 1.18, fill=BLUE)
    add_text(slide, "AI", 10.27, 2.05, 1.18, 1.18, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 9.80, 2.37, 10.23, 2.37, color=WHITE, width=2.2)
    add_arrow(slide, 10.25, 2.82, 9.84, 2.82, color=GOLD, width=2.2)
    add_box(slide, 9.18, 3.55, 1.75, 0.68, fill=WHITE, line=None)
    add_text(slide, "対 話", 9.18, 3.55, 1.75, 0.68, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 7.88, 3.58, 0.58, 0.42, fill=BLUE, rotation=218)
    add_chevron(slide, 10.98, 4.54, 0.58, 0.42, fill=GREEN, rotation=335)
    add_chevron(slide, 11.43, 1.16, 0.58, 0.42, fill=ORANGE, rotation=95)
    add_chip(slide, "提案", 7.55, 4.67, 1.17, fill=BLUE, color=WHITE)
    add_chip(slide, "検証", 9.47, 5.25, 1.17, fill=GREEN, color=WHITE)
    add_chip(slide, "判断", 11.25, 4.67, 1.17, fill=ORANGE, color=WHITE)
    add_text(
        slide,
        "Vive with Gemini",
        8.12,
        6.26,
        2.85,
        0.48,
        size=16,
        color="D9E8F5",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Put the public chapter link on the opening slide so attendees can scan it early.
    add_box(slide, 10.95, 5.25, 1.90, 2.00, fill=WHITE, line=BLUE, line_width=1.2)
    add_text(slide, "第5章を公開", 11.12, 5.32, 1.56, 0.27, size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    pic = slide.shapes.add_picture(
        str(qr_path), Inches(11.18), Inches(5.66), width=Inches(1.44), height=Inches(1.44)
    )
    try:
        pic.click_action.hyperlink.address = PUBLIC_URL
    except Exception:
        pass
    add_text(slide, "公開URL", 11.23, 7.08, 1.32, 0.14, size=8, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_note(
        slide,
        """
目安 0:00–0:50

「AIに仕事を丸ごと任せた成功談」ではありません。AIを相棒として、設計・実装・テスト・確認を一緒に進めた話です。
自己紹介は1文だけ。64歳から学び直し、2025年5月以降にAIとの対話で26アプリを試作し、現在は閉域MaaS基盤に取り組んでいます。
第5章の詳しい本文は右下のQRコードで公開していると、冒頭で短く案内します。

次へのつなぎ：相棒とは、判断までAIへ渡すことではありません。
""",
    )


def slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "WHY", "相棒とは、判断を委ねることではない", 2)
    add_chip(
        slide,
        "※ ROI＝価値↑・総コスト↓というプロジェクト内の“方言”",
        8.05,
        0.88,
        4.80,
        fill=ORANGE_LIGHT,
        color="A85E13",
        line=ORANGE,
        size=11,
    )

    add_card(
        slide,
        0.72,
        1.40,
        4.55,
        2.72,
        "人が持つもの",
        "目的・業務価値\n安全性・境界\n採用／不採用\n最終判断と責任",
        accent=ORANGE,
        fill=ORANGE_LIGHT,
        title_size=21,
        body_size=17,
        center=True,
    )
    add_circle(slide, 2.42, 1.08, 0.82, fill=ORANGE)
    add_text(slide, "人", 2.42, 1.08, 0.82, 0.82, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_card(
        slide,
        8.05,
        1.40,
        4.55,
        2.72,
        "AIに任せるもの",
        "選択肢の整理\n抜け漏れの指摘\n雛形・たたき台\n集計・確認支援",
        accent=BLUE,
        fill=BLUE_LIGHT,
        title_size=21,
        body_size=17,
        center=True,
    )
    add_circle(slide, 9.75, 1.08, 0.82, fill=BLUE)
    add_text(slide, "AI", 9.75, 1.08, 0.82, 0.82, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_arrow(slide, 5.48, 2.17, 7.85, 2.17, color=BLUE, width=3.0)
    add_arrow(slide, 7.85, 3.20, 5.48, 3.20, color=ORANGE, width=3.0)
    add_text(slide, "短い対話", 5.58, 2.38, 2.16, 0.55, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_text(
        slide,
        "チームの価値 ＝ 単位時間あたりのフィードバック量 × 質",
        1.03,
        4.47,
        11.25,
        0.58,
        size=23,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    loop_y = 5.32
    steps = [
        (1.45, "小さく\n試す", BLUE),
        (5.31, "早く\n確かめる", GREEN),
        (9.17, "次を\n決める", ORANGE),
    ]
    for x, label, color in steps:
        add_circle(slide, x, loop_y, 1.05, fill=color)
        add_text(slide, label, x, loop_y, 1.05, 1.05, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.63, 5.83, 5.17, 5.83, color=NAVY, width=2.6)
    add_arrow(slide, 6.49, 5.83, 9.03, 5.83, color=NAVY, width=2.6)
    add_arrow(slide, 10.32, 6.38, 2.05, 6.38, color="8AA3B8", width=1.8, dashed=True)

    add_footer(slide, "KEY", "短い対話の回数と質を上げ、価値を高めながら総コストを抑える")
    add_note(
        slide,
        """
目安 0:50–2:20

AIは選択肢、抜け漏れ、雛形、集計を支援します。目的・価値・安全性・採否・責任は人が持ちます。
完璧な質問を作り込む前に、小さく相談し、早く確かめて、次を決めます。
ここでのROI最大化は厳密な財務式ではなく、提供価値を高めながら手作業・待ち・再作業を含む総コストを下げる、Vive with Gemini固有の表現です。

次へのつなぎ：ただし、AIは確率的です。相棒として働ける枠組みが要ります。
""",
    )


def slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "HOW", "AIの強み①：決めた型を、毎回ぶれずに適用する", 3)

    # 1. What the team wants to build.
    add_card(
        slide,
        0.65,
        1.08,
        3.30,
        1.88,
        "作りたい仕様",
        "目的・要件・受入条件",
        accent=BLUE,
        fill=BLUE_LIGHT,
        title_size=20,
        body_size=15,
        number="1",
        center=True,
    )

    # 2. A template is one input, with local and global rules inside it.
    add_box(slide, 4.32, 1.08, 8.35, 1.88, fill=WHITE, line=PURPLE, line_width=1.5)
    add_box(slide, 4.32, 1.08, 0.10, 1.88, fill=PURPLE, line=None)
    add_circle(slide, 4.56, 1.26, 0.48, fill=PURPLE)
    add_text(slide, "2", 4.56, 1.26, 0.48, 0.48, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "テンプレート", 5.18, 1.18, 2.35, 0.56, size=20, color=INK, bold=True)
    add_text(
        slide,
        "判断に使う2種類の規約",
        8.00,
        1.22,
        4.15,
        0.46,
        size=13,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_box(slide, 4.72, 1.82, 3.55, 0.86, fill=PURPLE_LIGHT, line=PURPLE, line_width=1.1)
    add_text(slide, "ローカル規約", 4.88, 1.87, 1.55, 0.32, size=15, color=PURPLE, bold=True)
    add_text(
        slide,
        "設計・命名・実装・テスト",
        4.88,
        2.18,
        3.18,
        0.34,
        size=12,
        color=INK,
    )
    add_box(slide, 8.60, 1.82, 3.67, 0.86, fill=GREEN_LIGHT, line=GREEN, line_width=1.1)
    add_text(slide, "グローバル規約", 8.76, 1.87, 1.85, 0.32, size=15, color=GREEN, bold=True)
    add_text(
        slide,
        "共通標準・BP・設計原則",
        8.76,
        2.18,
        3.25,
        0.34,
        size=12,
        color=INK,
    )

    add_text(slide, "＋", 3.95, 1.61, 0.38, 0.58, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.30, 2.99, 4.10, 3.25, color=BLUE, width=2.4)
    add_arrow(slide, 8.50, 2.99, 8.05, 3.25, color=PURPLE, width=2.4)

    add_box(slide, 2.00, 3.28, 9.33, 0.86, fill=NAVY, line=None)
    add_circle(slide, 2.30, 3.42, 0.58, fill=WHITE)
    add_text(slide, "AI", 2.30, 3.42, 0.58, 0.58, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "決められた型を、毎回ぶれずに適用",
        3.05,
        3.34,
        7.85,
        0.72,
        size=21,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_arrow(slide, 6.67, 4.16, 6.67, 4.51, color=GREEN, width=3.0)

    # 3. Applying the agreed pattern improves both dimensions at once.
    add_box(slide, 0.92, 4.57, 11.49, 1.72, fill=GREEN_LIGHT, line=GREEN, line_width=1.8)
    add_circle(slide, 1.23, 4.82, 0.56, fill=GREEN)
    add_text(slide, "3", 1.23, 4.82, 0.56, 0.56, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "高品質 × 高生産性の成果物",
        2.02,
        4.67,
        9.45,
        0.61,
        size=24,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_chip(
        slide,
        "品質｜規約漏れ・ばらつきを抑える",
        2.04,
        5.46,
        4.30,
        0.46,
        fill=WHITE,
        color=GREEN,
        line=GREEN,
        size=13,
    )
    add_chip(
        slide,
        "生産性｜迷い・手戻りを減らす",
        6.97,
        5.46,
        4.30,
        0.46,
        fill=WHITE,
        color=BLUE,
        line=BLUE,
        size=13,
    )

    add_footer(slide, "KEY", "仕様と2つの規約を渡すと、AIの確実な適用が品質・生産性に直結する")
    add_note(
        slide,
        """
目安 2:20–3:50

AIには2つの側面があります。ここで強調するのは、パターン化した決め事を、毎回ぶれずに適用できる側面です。
作りたい仕様に加え、テンプレートとしてローカル規約とグローバル規約を渡します。
すると、規約漏れやばらつきを抑えて品質を上げながら、迷いや手戻りを減らし、生産性も上げられます。
チェックリストやテスト、最終判断は引き続き人が持ちますが、この図では「決めた型の確実な適用」が品質と生産性に直結する点に絞ります。

次へのつなぎ：その生産性を、顧客から学ぶ速さへ変えていきます。
""",
    )


def slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "HOW", "早く動くものを見せる。その反応が、価値になる", 4)

    # Left: a working slice that the customer can actually touch.
    add_box(slide, 0.65, 1.18, 3.55, 3.62, fill=BLUE_LIGHT, line=BLUE, line_width=1.6)
    add_box(slide, 0.65, 1.18, 0.10, 3.62, fill=BLUE, line=None)
    add_text(slide, "小さくても動くもの", 0.92, 1.36, 3.00, 0.54, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.15, 2.06, 2.55, 1.34, fill=WHITE, line=BLUE, radius=False, line_width=2.0)
    add_box(slide, 1.34, 2.23, 2.17, 0.18, fill=BLUE, line=None, radius=False)
    add_box(slide, 1.34, 2.55, 0.50, 0.58, fill=GRAY_LIGHT, line=None, radius=False)
    add_box(slide, 2.01, 2.55, 1.50, 0.20, fill=GREEN_LIGHT, line=GREEN, radius=False, line_width=0.8)
    add_box(slide, 2.01, 2.91, 1.04, 0.20, fill=ORANGE_LIGHT, line=ORANGE, radius=False, line_width=0.8)
    add_box(slide, 2.26, 3.40, 0.35, 0.30, fill=BLUE, line=None, radius=False)
    add_box(slide, 1.84, 3.68, 1.18, 0.12, fill=BLUE, line=None, radius=False)
    add_chip(
        slide,
        "AI＋自動化で早くつくる",
        1.02,
        4.15,
        2.82,
        0.42,
        fill=BLUE,
        color=WHITE,
        size=12,
    )

    # Right: concrete reactions after the customer touches the working slice.
    add_box(slide, 9.13, 1.18, 3.55, 3.62, fill=ORANGE_LIGHT, line=ORANGE, line_width=1.6)
    add_box(slide, 9.13, 1.18, 0.10, 3.62, fill=ORANGE, line=None)
    add_text(slide, "お客さまの反応", 9.42, 1.36, 2.98, 0.54, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_circle(slide, 9.52, 2.06, 0.54, fill=ORANGE)
    add_box(slide, 9.43, 2.61, 0.72, 0.78, fill=ORANGE, line=None)
    add_box(slide, 10.34, 2.00, 1.92, 0.54, fill=WHITE, line=GREEN, line_width=1.0)
    add_text(slide, "使える！", 10.39, 2.00, 1.82, 0.54, size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 10.34, 2.71, 1.92, 0.54, fill=WHITE, line=ORANGE, line_width=1.0)
    add_text(slide, "ここが使いにくい", 10.39, 2.71, 1.82, 0.54, size=12, color="A85E13", bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 10.34, 3.42, 1.92, 0.54, fill=WHITE, line=PURPLE, line_width=1.0)
    add_text(slide, "次はこれが欲しい", 10.39, 3.42, 1.82, 0.54, size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_chip(
        slide,
        "触った反応が返る",
        9.50,
        4.15,
        2.82,
        0.42,
        fill=ORANGE,
        color=WHITE,
        size=12,
    )

    # The two arrows form one conversation, not a staged process.
    add_text(slide, "早く見せる", 4.64, 1.54, 4.06, 0.42, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.34, 2.18, 8.94, 2.18, color=BLUE, width=3.0)
    add_arrow(slide, 8.94, 3.62, 4.34, 3.62, color=ORANGE, width=3.0)
    add_text(slide, "生のフィードバック", 4.64, 3.78, 4.06, 0.42, size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_circle(slide, 6.12, 2.46, 1.10, fill=NAVY)
    add_text(slide, "短い\n往復", 6.12, 2.46, 1.10, 1.10, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.67, 4.19, 6.67, 4.73, color=GREEN, width=3.0)

    # The loop itself creates learning and therefore customer value.
    add_box(slide, 0.92, 4.82, 11.49, 1.55, fill=GREEN_LIGHT, line=GREEN, line_width=1.8)
    add_circle(slide, 1.23, 5.22, 0.68, fill=GREEN)
    add_text(slide, "↑", 1.23, 5.22, 0.68, 0.68, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "この往復そのものが、顧客価値を育てる",
        2.04,
        4.94,
        9.42,
        0.60,
        size=23,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    value_chips = [
        (2.12, "ズレが見える"),
        (5.16, "次を決められる"),
        (8.20, "価値を早く届ける"),
    ]
    for x, label in value_chips:
        add_chip(slide, label, x, 5.72, 2.58, 0.42, fill=WHITE, color=GREEN, line=GREEN, size=12)

    add_footer(slide, "KEY", "早く動かし、早く反応を得る。その往復が、そのまま顧客価値になる")
    add_note(
        slide,
        """
目安 3:50–5:40

ここは工程を順番に進める話ではありません。小さくても動くものを早くつくり、お客さまに触ってもらいます。
そうすると、「使える」「ここが使いにくい」「次はこれが欲しい」という生のフィードバックが返ります。
価値は完成時に初めて生まれるのではありません。この往復でズレを早く見つけ、次に作るべきものを確かめること自体が、すでに顧客価値へつながっています。
AIと自動化は、この往復を速める手段です。改善はAI単独の効果ではなく、作業分割やスクラム運営を含むチーム全体の結果です。

次へのつなぎ：同じ考え方を、現在の閉域MaaS活動へ持ち込んでいます。
""",
    )


def slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "NOW", "現在の閉域MaaS活動は、4本柱", 5)

    pillars = [
        (0.48, "1", "AI協働開発", "Copilot / Codex\n境界を守る", BLUE, BLUE_LIGHT),
        (3.68, "2", "MaaS実行基盤", "Kong / OpenShift\nKubernetes / GPU", ORANGE, ORANGE_LIGHT),
        (6.88, "3", "性能評価自動化", "GuideLLM / Excel\nCopilot Agent", GREEN, GREEN_LIGHT),
        (10.08, "4", "適応型開発", "GitLab Issue\n1週間スプリント", PURPLE, PURPLE_LIGHT),
    ]
    for x, num, title, body, accent, fill in pillars:
        add_card(slide, x, 1.10, 2.77, 2.02, title, body, accent=accent, fill=fill, title_size=17, body_size=14, number=num, center=True)
        add_arrow(slide, x + 1.38, 3.16, 6.67, 3.60, color=accent, width=1.8)
    add_box(slide, 4.44, 3.45, 4.45, 0.72, fill=NAVY, line=None)
    add_text(slide, "会話駆動型MaaSの価値提供ループ", 4.44, 3.45, 4.45, 0.72, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    status = [
        (0.55, "構築済み", "GuideLLM共通手順\n評価ブックAgent", GREEN, GREEN_LIGHT),
        (4.50, "実施中", "MaaS・GPU検証\n1週間スプリント", BLUE, BLUE_LIGHT),
        (8.45, "今 後", "CI/CD・監視\n一気通貫の自動化", PURPLE, PURPLE_LIGHT),
    ]
    for x, title, body, accent, fill in status:
        add_card(slide, x, 4.52, 3.84, 1.58, title, body, accent=accent, fill=fill, title_size=17, body_size=14, center=True)
    add_chip(slide, "効果は測定中", 5.35, 6.20, 2.63, fill=RED_LIGHT, color=RED, line=RED, size=13)

    add_footer(slide, "KEY", "構築済み・実施中・今後を分け、完成形や効果を先に決めない")
    add_note(
        slide,
        """
目安 5:40–7:20

顧客管理刷新とは別に、現在は閉域向けの社内MaaS基盤に取り組んでいます。
4本柱は、AI協働開発、MaaS実行基盤、性能評価自動化、GitLab Issueと1週間スプリントです。
共通ベンチ手順と評価ブックAgentは構築済み。MaaS、GPU検証、1週間スプリントは実施中。CI/CD、監視、一気通貫自動化は今後です。削減時間・削減率などの効果は測定中です。

次へのつなぎ：閉域で最初に設計したのは、AI利用の境界です。
""",
    )


def slide_6(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "SECURITY", "AIを使える領域と、守る領域を分ける", 6)

    add_card(
        slide,
        0.55,
        1.18,
        3.48,
        4.48,
        "承認済みAI領域",
        "Microsoft 365 Copilot\n会議・要望・Excel分析\n\nVS Code上のCodex\n設計・実装・テスト支援\n\n入力が許可された情報のみ",
        accent=BLUE,
        fill=BLUE_LIGHT,
        title_size=20,
        body_size=15,
        center=True,
    )
    add_card(
        slide,
        4.92,
        1.55,
        3.48,
        3.72,
        "統制された受け渡し",
        "1  情報と成果物を確認\n\n2  人が内容・安全性をレビュー\n\n3  承認済み成果物だけを反映",
        accent=ORANGE,
        fill=ORANGE_LIGHT,
        title_size=20,
        body_size=15,
        center=True,
    )
    add_card(
        slide,
        9.28,
        1.18,
        3.48,
        4.48,
        "エアギャップ／閉域",
        "Kong Gateway\n↓\nOpenShift / Kubernetes\n↓\nGPU\n\nMaaSを配置・実行・計測",
        accent=GREEN,
        fill=GREEN_LIGHT,
        title_size=20,
        body_size=15,
        center=True,
    )
    add_arrow(slide, 4.12, 3.30, 4.77, 3.30, color=BLUE, width=3.0)
    add_text(slide, "成果物", 4.10, 2.78, 0.70, 0.35, size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 8.49, 3.30, 9.14, 3.30, color=GREEN, width=3.0)
    add_text(slide, "所定手順", 8.43, 2.78, 0.84, 0.35, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 1.10, 6.02, 11.12, 0.55, fill=RED_LIGHT, line=RED, line_width=1.2)
    add_text(
        slide,
        "入力が許可されていない情報・認証情報  ── × ──▶  外部AI",
        1.27,
        6.03,
        10.78,
        0.52,
        size=16,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "KEY", "クラウドAIから閉域MaaSへ直接接続せず、人の承認を境界に置く")
    add_note(
        slide,
        """
目安 7:20–9:00

外部接続が許可された承認済み領域だけでCopilotとCodexを使い、入力が許可された情報だけを渡します。
人が成果物の内容と安全性をレビューし、承認済み成果物だけを所定手順で閉域へ反映します。
閉域MaaSからクラウドAIへ直接接続する構成ではありません。閉域側ではKongを入口に、OpenShift／KubernetesがPodとGPUを管理します。

次へのつなぎ：閉域MaaSでモデルを提供するには、性能を同じ条件で比べる必要があります。
""",
    )


def slide_7(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "BENCHMARK", "2つの入力軸を、同じ時間軸で比べる", 7)

    add_card(slide, 0.45, 1.08, 3.05, 2.10, "入力1｜GuideLLM負荷", "Request rate\nConcurrency\n入出力Token / 試験時間", accent=BLUE, fill=BLUE_LIGHT, title_size=16, body_size=14, center=True)
    add_card(slide, 0.45, 3.48, 3.05, 2.10, "入力2｜vLLM設定", "max-model-len / seqs\nbatched-tokens\nGPU memory設定値", accent=PURPLE, fill=PURPLE_LIGHT, title_size=16, body_size=13, center=True)

    add_box(slide, 4.35, 2.28, 2.55, 2.08, fill=GREEN_LIGHT, line=GREEN, line_width=1.8)
    add_text(slide, "OpenShift上の\nvLLM / MaaS", 4.53, 2.52, 2.19, 0.86, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "同じモデル・GPU構成・版", 4.52, 3.55, 2.20, 0.36, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.60, 2.15, 4.20, 2.82, color=BLUE, width=2.6)
    add_arrow(slide, 3.60, 4.48, 4.20, 3.84, color=PURPLE, width=2.6)

    add_card(slide, 7.68, 1.08, 3.17, 2.10, "出力1｜GuideLLM性能", "Throughput / Latency\nTTFT / 成功・エラー件数", accent=GREEN, fill=GREEN_LIGHT, title_size=16, body_size=14, center=True)
    add_card(slide, 7.68, 3.48, 3.17, 2.10, "出力2｜基盤メトリクス", "Queue / Wait / KV\nCPU / GPU / Tensor Core", accent=NAVY, fill=BLUE_LIGHT, title_size=16, body_size=14, center=True)
    add_arrow(slide, 7.02, 2.82, 7.54, 2.14, color=GREEN, width=2.6)
    add_arrow(slide, 7.02, 3.84, 7.54, 4.50, color=NAVY, width=2.6)

    add_box(slide, 11.33, 2.23, 1.55, 2.20, fill=ORANGE_LIGHT, line=ORANGE, line_width=1.6)
    add_text(slide, "人が\n次の条件を\n判断", 11.48, 2.54, 1.25, 1.28, size=17, color="A85E13", bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.93, 2.18, 11.20, 2.72, color=ORANGE, width=2.2)
    add_arrow(slide, 10.93, 4.50, 11.20, 3.95, color=ORANGE, width=2.2)
    add_arrow(slide, 12.10, 4.60, 12.10, 5.82, color="8AA3B8", width=1.6, dashed=True)
    add_arrow(slide, 12.10, 5.82, 1.98, 5.82, color="8AA3B8", width=1.6, dashed=True)

    add_chip(slide, "性能評価 ≠ 回答品質評価", 0.63, 6.05, 3.35, fill=RED_LIGHT, color=RED, line=RED, size=12)
    add_chip(slide, "時間一致 ≠ 因果", 4.99, 6.05, 2.40, fill=GRAY_LIGHT, color=MUTED, line=LINE, size=12)
    add_chip(slide, "自動チューニングではない", 8.37, 6.05, 3.35, fill=GRAY_LIGHT, color=MUTED, line=LINE, size=12)

    add_footer(slide, "KEY", "外形性能と内部状態を突き合わせ、次の検証条件を人が決める")
    add_note(
        slide,
        """
目安 9:00–11:30

入力軸は2つです。GuideLLM側でリクエストレート、同時実行数、トークン数などの負荷を変えます。vLLM側では最大コンテキスト長、シーケンス数、バッチトークン、GPUメモリ設定などを変えます。
GuideLLMで利用者視点のスループット、レイテンシー、TTFTなどを見て、同じ時間帯のPromQLでQueue、CPU、GPUなどの内部状態を見ます。
GuideLLMは性能評価であり回答品質評価ではありません。時間が一致しても因果は断定しません。自動チューニングではなく、人が次の条件を決めます。

次へのつなぎ：ここで最も手作業が多かったのが、結果を評価ブックへまとめる部分です。
""",
    )


def slide_8(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "SHOWCASE", "3つの入力から、共有できる評価ブックへ", 8)
    add_chip(slide, "構築済み｜プロジェクト関係者へ公開済み", 8.58, 0.87, 4.20, fill=GREEN_LIGHT, color=GREEN, line=GREEN, size=12)

    # Input column
    add_box(slide, 0.48, 1.28, 2.55, 4.82, fill=BLUE_LIGHT, line=BLUE, line_width=1.6)
    add_chip(slide, "INPUT｜3種類", 0.75, 1.54, 2.01, fill=BLUE, color=WHITE, size=14)
    for i, label in enumerate(["Excelテンプレート", "GuideLLM結果\ntarファイル", "Pod起動ログ"]):
        y = 2.22 + i * 1.08
        add_box(slide, 0.76, y, 1.99, 0.82, fill=WHITE, line="9DC6F8", line_width=1.1)
        add_text(slide, label, 0.84, y + 0.05, 1.83, 0.72, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "許可された評価データのみ", 0.77, 5.57, 1.97, 0.30, size=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Agent column
    add_box(slide, 3.74, 1.28, 2.78, 4.82, fill=PURPLE_LIGHT, line=PURPLE, line_width=1.6)
    add_chip(slide, "AGENT｜作成支援", 4.05, 1.54, 2.16, fill=PURPLE, color=WHITE, size=14)
    add_circle(slide, 4.76, 2.33, 0.74, fill=PURPLE)
    add_text(slide, "A", 4.76, 2.33, 0.74, 0.74, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Microsoft 365\nCopilot Agent", 4.03, 3.17, 2.20, 0.90, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 4.05, 4.36, 2.16, 0.78, fill=WHITE, line="C8B8DF")
    add_text(slide, "ベンチ評価ブック作成", 4.15, 4.43, 1.96, 0.64, size=14, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # Output column
    add_box(slide, 7.23, 1.28, 2.78, 4.82, fill=GREEN_LIGHT, line=GREEN, line_width=1.6)
    add_chip(slide, "OUTPUT｜標準化", 7.52, 1.54, 2.20, fill=GREEN, color=WHITE, size=14)
    add_box(slide, 7.65, 2.36, 1.94, 0.76, fill=WHITE, line="A8D7BC")
    add_text(slide, "評価ブック", 7.74, 2.40, 1.76, 0.66, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["比較表", "グラフ", "所定の列・形式"]):
        y = 3.42 + i * 0.67
        add_chip(slide, label, 7.70, y, 1.84, 0.46, fill=WHITE, color=INK, line="A8D7BC", size=13)

    # Human column
    add_box(slide, 10.71, 1.28, 2.16, 4.82, fill=ORANGE_LIGHT, line=ORANGE, line_width=1.6)
    add_chip(slide, "HUMAN", 10.96, 1.54, 1.66, fill=ORANGE, color=WHITE, size=14)
    add_circle(slide, 11.40, 2.29, 0.78, fill=ORANGE)
    add_text(slide, "人", 11.40, 2.29, 0.78, 0.78, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "元データ確認\n\nモデル選定\n\nGPU構成判断", 10.96, 3.18, 1.66, 2.10, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)

    add_arrow(slide, 3.10, 3.66, 3.61, 3.66, color=BLUE, width=2.8)
    add_arrow(slide, 6.61, 3.66, 7.10, 3.66, color=PURPLE, width=2.8)
    add_arrow(slide, 10.10, 3.66, 10.58, 3.66, color=ORANGE, width=2.8)
    add_chip(slide, "入力 → Agent → 評価ブック → 人の判断", 3.95, 6.20, 5.44, fill=NAVY, color=WHITE, size=14)

    add_footer(slide, "KEY", "作成はAgent。元データの正しさ、モデル採用、GPU構成の判断は人")
    add_note(
        slide,
        """
目安 11:30–14:30（最大の見せ場）

独自Agent「ベンチ評価ブック作成」は構築済みで、プロジェクト関係者へ公開しています。
入力はExcelテンプレート、GuideLLM結果のtarファイル、Pod起動ログの3つです。Agentが所定形式の評価ブック、比較表、グラフを作ります。
渡すのは社内ルール上許可された評価データだけです。元データ、集計結果、表・グラフを人が確認し、モデル採用とGPU構成は人が最終判断します。Copilotが評価の正しさを保証するわけではありません。
効率化にはつながりましたが、削減時間・削減率の定量評価は継続中です。

次へのつなぎ：作った評価を一度で終わらせず、次のGitLab Issueへ戻します。
""",
    )


def slide_9(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "LOOP", "1週間の中で、AIと「分単位」に対話する", 9)

    loop_x = 0.42
    loop_y = 1.10
    loop_w = 8.58
    loop_h = 5.40
    add_box(slide, loop_x, loop_y, loop_w, loop_h, fill=WHITE, line=LINE, line_width=1.2)
    add_chip(slide, "外側：1週間スプリント", 2.82, 1.23, 3.76, fill=NAVY, color=WHITE, size=14)
    nodes = [
        (0.72, 1.82, "ユーザー\n対話", BLUE),
        (3.25, 1.82, "GitLab\nIssue", PURPLE),
        (5.78, 1.82, "AI＋人で\n開発", BLUE),
        (5.78, 4.98, "人が\n承認", ORANGE),
        (3.25, 4.98, "閉域MaaS\n実行", GREEN),
        (0.72, 4.98, "GuideLLM\n計測", GREEN),
    ]
    for x, y, label, color in nodes:
        add_box(slide, x, y, 1.82, 0.92, fill=WHITE, line=color, line_width=1.6)
        add_text(slide, label, x + 0.08, y + 0.05, 1.66, 0.82, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.62, 2.28, 3.12, 2.28, color=NAVY, width=2.2)
    add_arrow(slide, 5.15, 2.28, 5.65, 2.28, color=NAVY, width=2.2)
    add_arrow(slide, 6.70, 2.85, 6.70, 4.84, color=NAVY, width=2.2)
    add_arrow(slide, 5.65, 5.44, 5.15, 5.44, color=NAVY, width=2.2)
    add_arrow(slide, 3.12, 5.44, 2.62, 5.44, color=NAVY, width=2.2)
    add_arrow(slide, 1.63, 4.84, 1.63, 2.85, color=ORANGE, width=2.2, dashed=True)

    # Inner cadence: minute-by-minute human/AI dialogue inside the weekly sprint.
    add_box(slide, 2.15, 2.85, 4.90, 1.55, fill=BLUE_LIGHT, line=BLUE, line_width=1.8)
    add_chip(slide, "内側：分単位の対話", 3.36, 2.98, 2.50, fill=BLUE, color=WHITE, size=14)
    add_circle(slide, 2.52, 3.35, 0.72, fill=ORANGE)
    add_text(slide, "人", 2.52, 3.35, 0.72, 0.72, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_circle(slide, 5.97, 3.35, 0.72, fill=BLUE)
    add_text(slide, "AI", 5.97, 3.35, 0.72, 0.72, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.39, 3.52, 5.82, 3.52, color=BLUE, width=2.5)
    add_arrow(slide, 5.82, 3.91, 3.39, 3.91, color=ORANGE, width=2.5)
    add_text(slide, "質問 → 提案 → 確認 → 次の質問", 3.39, 3.50, 2.43, 0.42, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "考え込む時間を、相談と検証の回数へ変える", 2.38, 4.48, 4.44, 0.30, size=12, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 9.35, 1.10, 3.55, 5.40, fill=BLUE_LIGHT, line=NAVY, line_width=1.3)
    add_text(slide, "価値を生むリズム", 9.70, 1.38, 2.85, 0.42, size=17, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, "AIとの会話回数 ↑", 9.72, 1.96, 2.80, 0.53, fill=BLUE, color=WHITE, size=17)
    add_text(slide, "↓", 10.78, 2.57, 0.70, 0.45, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, "学び・価値 ↑", 9.72, 3.07, 2.80, 0.53, fill=GREEN, color=WHITE, size=17)
    add_rich_text(
        slide,
        [("価値 ∝ ", 19, NAVY, True), ("対話回数", 19, BLUE, True), (" × 質", 19, ORANGE, True)],
        9.55,
        3.78,
        3.15,
        0.62,
        align=PP_ALIGN.CENTER,
    )
    add_box(slide, 9.67, 4.60, 2.92, 0.92, fill=WHITE, line=ORANGE, line_width=1.2)
    add_text(slide, "会話の質を守るもの\n規約・テスト・人の判断", 9.80, 4.69, 2.66, 0.72, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_chip(slide, "採否と責任は人", 9.91, 5.85, 2.42, fill=ORANGE_LIGHT, color="A85E13", line=ORANGE, size=13)

    add_footer(slide, "KEY", "AIと分単位に会話する。対話回数に比例して、学びと価値を積み上げる")
    add_note(
        slide,
        """
目安 14:30–16:30

外側は1週間スプリントです。ユーザー要望をGitLab Issueへ整理し、毎週、検証結果または動く成果物を見せてフィードバックを得ます。
しかし、価値を生む中心は週1回の会話だけではありません。スプリントの内側で、人がAIへ分単位で質問し、AIの提案を確認し、すぐ次の質問へ進みます。
考え込んで止まる時間を相談と検証の回数へ変え、AIとの対話回数に比例して学びと価値を積み上げる、というのが最も伝えたい点です。
ただし回数だけを増やすのではなく、許可情報、規約、自動テスト、人の判断によって会話の質を守り、採否と責任は人が持ちます。

次へのつなぎ：最後に、成功談ではなく再現できる「任せ方」を3つ残します。
""",
    )


def slide_10(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "TIPS", "成功談より持ち帰ってほしい「任せ方の勘所」", 10)

    tips = [
        ("1", "小さく任せる", "候補・雛形・抜け漏れ確認へ切り出す", BLUE, BLUE_LIGHT),
        ("2", "判断基準を先に渡す", "許可情報・規約・受入条件・出力形式", PURPLE, PURPLE_LIGHT),
        ("3", "証拠で受け取る", "差分・テスト・元データを人が確認する", GREEN, GREEN_LIGHT),
    ]
    for i, (num, title, body, accent, fill) in enumerate(tips):
        y = 1.10 + i * 1.36
        add_card(slide, 0.55, y, 12.23, 1.12, title, body, accent=accent, fill=fill, title_size=20, body_size=14, number=num)

    add_box(slide, 0.75, 5.32, 11.83, 0.62, fill=ORANGE_LIGHT, line=ORANGE, line_width=1.2)
    add_text(slide, "AIの誤り → 修正 → 規約・テスト・次のIssueへ戻す", 0.95, 5.35, 11.43, 0.56, size=18, color="A85E13", bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 0.55, 6.12, 12.23, 0.63, fill=NAVY, line=None)
    add_rich_text(
        slide,
        [
            ("AIに任せる：", 16, WHITE, True),
            ("作業と提案　", 16, GOLD, True),
            ("人が手放さない：", 16, WHITE, True),
            ("目的・境界・判断・責任", 16, ORANGE, True),
        ],
        0.78,
        6.16,
        11.77,
        0.54,
        align=PP_ALIGN.CENTER,
    )

    add_note(
        slide,
        f"""
目安 16:30–18:00（約2分余白）

成功談で終わらず、任せ方を3つ持ち帰ってください。
1. 小さく任せる。候補、雛形、抜け漏れ確認に切り出す。
2. 判断基準を先に渡す。許可情報、規約、受入条件、出力形式を共有する。
3. 証拠で受け取る。差分、テスト、元データを人が確認する。
AIの誤りは隠さず、修正し、規約・テスト・次のIssueへ戻します。

締め：AIに任せるのは作業と提案。人が手放さないのは目的・境界・判断・責任です。
詳細、前提、構築済み／実施中／今後の区別は、第1ページのQRコードから見られる第5章で公開しています：{PUBLIC_URL}
""",
    )


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    qr = qrcode.QRCode(version=None, box_size=10, border=2)
    qr.add_data(PUBLIC_URL)
    qr.make(fit=True)
    qr.make_image(fill_color=f"#{NAVY_DARK}", back_color=f"#{WHITE}").save(QR_FILE)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "AIは「道具」でなく「相棒」"
    prs.core_properties.subject = "AI開発事例共有会・10枚版"
    prs.core_properties.author = "原 桂介"
    prs.core_properties.keywords = "AI駆動開発, アジャイル, MaaS, GuideLLM, Copilot, Codex"
    prs.core_properties.comments = "Generated from chapter 5 through 5.3 of Vive with Gemini."

    slide_1(prs, QR_FILE)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)

    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    path = build()
    print(path)
    print(f"slides=10 url={PUBLIC_URL}")
