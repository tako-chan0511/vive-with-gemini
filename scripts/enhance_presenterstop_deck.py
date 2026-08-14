#!/usr/bin/env python3
"""Add public references and a final takeaway slide to the 24-slide deck.

The 24-slide PresenterStop deck does not have a source generator in this
repository.  This script therefore treats that deck as an input artifact and
creates a new 25-slide version without overwriting the original.

Run with:

    PYTHONPATH=/tmp/vive_ppt_deps \
      python3 scripts/enhance_presenterstop_deck.py
"""

from __future__ import annotations

import unicodedata
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import quote

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
INPUT_FILE = (
    PRESENTATION_DIR
    / "ai-developer-meeting-vive-with-gemini-presenterstop-24slides.pptx"
)
OUTPUT_FILE = (
    PRESENTATION_DIR
    / "ai-developer-meeting-vive-with-gemini-presenterstop-25slides.pptx"
)
QR_FILE = Path("/tmp/vive-presenterstop-final-qr.png")

BASE = "https://tako-chan0511.github.io/vive-with-gemini/"
PAGE_5 = f"{BASE}ai-agile-vive-with-gemini-5.html"
PAGE_51 = f"{BASE}ai-agile-vive-with-gemini-5-1.html"
PAGE_52 = f"{BASE}ai-agile-vive-with-gemini-5-2.html"
PAGE_53 = f"{BASE}ai-agile-vive-with-gemini-5-3.html"

SCRUM_GUIDE = (
    "https://scrumguides.org/docs/scrumguide/v2020/"
    "2020-Scrum-Guide-Japanese.pdf"
)
NIST_AI_RMF = "https://airc.nist.gov/airmf-resources/airmf/5-sec-core/"
M365_COPILOT_SECURITY = (
    "https://learn.microsoft.com/en-us/microsoft-365/copilot/"
    "security-microsoft-365-copilot"
)
CODEX_SECURITY = "https://learn.chatgpt.com/docs/agent-approvals-security"
KONG_GATEWAY = "https://developer.konghq.com/gateway/"
GUIDELLM = "https://github.com/vllm-project/guidellm"
VLLM_METRICS = "https://docs.vllm.ai/en/latest/design/metrics/"
GITLAB_ITERATIONS = "https://docs.gitlab.com/user/group/iterations/"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_JP = "Yu Gothic"

NAVY = "0B1F33"
NAVY_CARD = "15344F"
BLUE = "42A5F5"
BLUE_LINK = "64B5F6"
ORANGE = "FF9F43"
GREEN = "4CD97B"
PURPLE = "A78BFA"
WHITE = "FFFFFF"
TEXT_LIGHT = "DCEAF4"
MUTED = "AFC4D6"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def heading_url(page: str, heading_id: str | None = None) -> str:
    """Build a VitePress heading URL using the site's decomposed Unicode IDs."""

    if not heading_id:
        return page
    fragment = unicodedata.normalize("NFD", heading_id)
    return f"{page}#{quote(fragment, safe='-._~')}"


@dataclass(frozen=True)
class SlideReference:
    label: str
    page_url: str
    target_url: str
    official_urls: tuple[tuple[str, str], ...] = ()


def ref(
    label: str,
    page_url: str,
    heading_id: str | None = None,
    *official_urls: tuple[str, str],
) -> SlideReference:
    return SlideReference(
        label=label,
        page_url=page_url,
        target_url=heading_url(page_url, heading_id),
        official_urls=tuple(official_urls),
    )


REFERENCES: dict[int, SlideReference] = {
    1: ref(
        "発表の詳しい内容",
        PAGE_5,
        'aiは「道具」でなく「相棒」',
    ),
    2: ref("自己紹介・26アプリ", PAGE_5, "原-桂介-keisuke-hara"),
    3: ref("本日の流れ", PAGE_5, "本日の流れ"),
    4: ref(
        "なぜフィードバックループか",
        PAGE_51,
        "_1-なぜフィードバックループなのか",
    ),
    5: ref(
        "AIと人の役割分担",
        PAGE_51,
        '_2-「相棒」とは何か',
    ),
    6: ref(
        "自然発生アジャイル",
        PAGE_51,
        "自然発生アジャイル",
        ("スクラムガイド（日本語）", SCRUM_GUIDE),
    ),
    7: ref(
        "AI活用の前提",
        PAGE_52,
        "_1-実践の前提",
        ("NIST AI RMF Core", NIST_AI_RMF),
    ),
    8: ref(
        "AI活用の枠組み",
        PAGE_52,
        '_2-ai活用を支える「枠組み」',
    ),
    9: ref(
        "フィードバックを速くする",
        PAGE_52,
        "_4-aiと自動化でフィードバックを速くする",
    ),
    10: ref(
        "4本柱と現在地",
        PAGE_53,
        "now-閉域環境向けお客様社内maas-モデル提供サービス-基盤",
    ),
    11: ref(
        "セキュリティ境界",
        PAGE_53,
        "_1️⃣-ai利用領域と閉域領域を分離する",
        ("Microsoft 365 Copilot Security", M365_COPILOT_SECURITY),
        ("Codex Agent approvals & security", CODEX_SECURITY),
    ),
    12: ref(
        "基盤の責任分離",
        PAGE_53,
        "_3️⃣-kong、openshift-kubernetes、gpuの責任を分ける",
        ("Kong Gateway", KONG_GATEWAY),
    ),
    13: ref(
        "性能評価と品質評価",
        PAGE_53,
        "_4️⃣-多様なモデルを短時間で比較する性能ベンチマーク自動化",
        ("GuideLLM", GUIDELLM),
    ),
    14: ref(
        "2つの入力条件",
        PAGE_53,
        "入力1-guidellm側の負荷条件を変える",
        ("GuideLLM", GUIDELLM),
    ),
    15: ref(
        "共通手順での実行",
        PAGE_53,
        "権限を持つ社内利用者が同じ手順で実行する",
        ("GuideLLM", GUIDELLM),
    ),
    16: ref(
        "評価ブック作成Agent",
        PAGE_53,
        "評価データを保存し、copilot-agentで評価ブックを作成する",
    ),
    17: ref(
        "自動化フロー",
        PAGE_53,
        "maas性能ベンチマーク・分析の自動化フロー",
    ),
    18: ref(
        "外形性能と内部メトリクス",
        PAGE_53,
        "guidellm結果とvllm・基盤メトリクスの突き合わせ",
        ("vLLM Metrics", VLLM_METRICS),
    ),
    19: ref(
        "評価コストの計測",
        PAGE_53,
        "評価コストの低減を計測する",
    ),
    20: ref(
        "1週間スプリント",
        PAGE_53,
        "_5️⃣-gitlab-issueを利用した1週間スプリント",
        ("GitLab Iterations", GITLAB_ITERATIONS),
    ),
    21: ref(
        "4本柱のループ",
        PAGE_53,
        "_6️⃣-4本柱を1つのフィードバックループへ",
    ),
    22: ref(
        "ROI最大化の定義",
        PAGE_53,
        '_7️⃣-vive-with-geminiの-方言-「roi最大化」',
    ),
    23: ref("代表KPI", PAGE_53, "_8️⃣-代表kpi"),
    24: ref("結び", PAGE_53, "💬-結び"),
    25: ref("今回伝えたかったこと", PAGE_51, "_6-伝えたいこと"),
}


def set_background(slide, color: str = NAVY) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str | None = None,
    radius: bool = True,
    line_width: float = 1.4,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_takeaway_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    number: str,
    title: str,
    body: str,
    accent: str,
) -> None:
    add_box(slide, x, y, w, h, fill=NAVY_CARD, line=accent, line_width=1.5)
    add_box(slide, x, y, 0.10, h, fill=accent, radius=True)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x + 0.22),
        Inches(y + 0.20),
        Inches(0.48),
        Inches(0.48),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(accent)
    badge.line.fill.background()
    add_text(
        slide,
        number,
        x + 0.22,
        y + 0.20,
        0.48,
        0.48,
        size=15,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        title,
        x + 0.83,
        y + 0.10,
        w - 1.05,
        0.48,
        size=17,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.84,
        y + 0.55,
        w - 1.08,
        h - 0.62,
        size=12.5,
        color=TEXT_LIGHT,
        valign=MSO_ANCHOR.TOP,
    )


def add_reference_footer(slide, reference: SlideReference) -> None:
    """Add a compact, visible URL and link it to the matching page heading."""

    box = slide.shapes.add_textbox(
        Inches(3.45), Inches(6.965), Inches(6.78), Inches(0.29)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"公開HP ↗  {reference.page_url}"
    run.font.name = FONT_JP
    run.font.size = Pt(6.7)
    run.font.bold = True
    run.font.color.rgb = rgb(BLUE_LINK)
    run.font.underline = True
    run.hyperlink.address = reference.target_url
    box.click_action.hyperlink.address = reference.target_url


def append_reference_notes(slide, reference: SlideReference) -> None:
    frame = slide.notes_slide.notes_text_frame
    current = frame.text.rstrip()
    lines = [
        "",
        f"公開参照: {reference.label}",
        reference.target_url,
    ]
    if reference.official_urls:
        lines.append("関連する公式資料:")
        lines.extend(f"• {label}: {url}" for label, url in reference.official_urls)
    frame.text = current + "\n" + "\n".join(lines)


def add_final_slide(prs: Presentation, qr_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    add_box(slide, 0.55, 0.35, 1.42, 0.38, fill=BLUE)
    add_text(
        slide,
        "CLOSING",
        0.59,
        0.35,
        1.34,
        0.38,
        size=11,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "25",
        12.10,
        0.28,
        0.65,
        0.42,
        size=12,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        "今回、伝えたかったこと",
        0.62,
        0.72,
        12.05,
        0.78,
        size=40,
        bold=True,
    )
    add_text(
        slide,
        "人とAIの役割を分け、短いフィードバックループで価値を育てる。",
        0.70,
        1.48,
        11.95,
        0.42,
        size=17,
        color=MUTED,
        bold=True,
    )

    add_takeaway_card(
        slide,
        0.70,
        2.02,
        5.82,
        1.22,
        "1",
        "AIは「相棒」",
        "作業と提案はAIへ。\n最終判断と責任は人が持つ。",
        BLUE,
    )
    add_takeaway_card(
        slide,
        6.80,
        2.02,
        5.82,
        1.22,
        "2",
        "小さく、早く回す",
        "試す → 確かめる → 学ぶ →\n次の判断へ返す。",
        ORANGE,
    )
    add_takeaway_card(
        slide,
        0.70,
        3.43,
        5.82,
        1.22,
        "3",
        "境界と証拠で守る",
        "許可情報・規約・自動テスト・差分・\n人のレビューを組み合わせる。",
        GREEN,
    )
    add_takeaway_card(
        slide,
        6.80,
        3.43,
        5.82,
        1.22,
        "4",
        "学びを次へつなぐ",
        "人・AI・基盤・評価データをつなぎ、\n次のIssue／スプリントへ戻す。",
        PURPLE,
    )

    add_box(slide, 0.70, 4.90, 11.92, 1.42, fill=WHITE)
    add_text(
        slide,
        "AIに任せるのは、作業と提案。\n人が手放さないのは、目的・境界・判断・責任。",
        0.95,
        5.03,
        9.82,
        1.14,
        size=20,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    picture = slide.shapes.add_picture(
        str(qr_path), Inches(11.18), Inches(5.02), Inches(0.92), Inches(0.92)
    )
    picture.click_action.hyperlink.address = REFERENCES[25].target_url
    add_text(
        slide,
        "詳しい本文 ↗",
        10.95,
        5.93,
        1.38,
        0.24,
        size=7,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    ).click_action.hyperlink.address = REFERENCES[25].target_url

    add_text(
        slide,
        "PresenterStop 25  •  14:10",
        0.62,
        6.98,
        3.20,
        0.28,
        size=10,
        color=MUTED,
        bold=True,
    )
    add_text(
        slide,
        "Vive with Gemini",
        10.45,
        6.98,
        2.28,
        0.28,
        size=10,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_reference_footer(slide, REFERENCES[25])

    slide.notes_slide.notes_text_frame.text = (
        "PresenterStop 25/25｜目安 14:10\n"
        "停止点: 最後に持ち帰ってほしい行動原則\n"
        "話す要点:\n"
        "• 今日お伝えしたかったのは、AIに仕事を丸ごと渡すことではありません。\n"
        "• 作業と提案をAIへ任せ、目的・境界・判断・責任を人が持ちます。\n"
        "• 小さく試し、証拠で確かめ、学びを次のIssueへ返します。\n"
        "• この反復を人・AI・基盤・評価データで回し続けることが、"
        "Vive with Geminiの実践です。"
    )
    append_reference_notes(slide, REFERENCES[25])


def build() -> Path:
    if not INPUT_FILE.exists():
        raise SystemExit(f"Input deck not found: {INPUT_FILE}")

    prs = Presentation(INPUT_FILE)
    if len(prs.slides) != 24:
        raise SystemExit(f"Expected 24 input slides, found {len(prs.slides)}")
    if round(prs.slide_width / 914400, 3) != SLIDE_W:
        raise SystemExit("Unexpected slide width")
    if round(prs.slide_height / 914400, 3) != SLIDE_H:
        raise SystemExit("Unexpected slide height")

    for slide_number, slide in enumerate(prs.slides, start=1):
        reference = REFERENCES[slide_number]
        # Slide 1 already has a prominent visible URL and QR code.
        if slide_number != 1:
            add_reference_footer(slide, reference)
        append_reference_notes(slide, reference)

    qr = qrcode.QRCode(version=4, box_size=8, border=2)
    qr.add_data(REFERENCES[25].target_url)
    qr.make(fit=True)
    qr.make_image(fill_color="black", back_color="white").save(QR_FILE)

    add_final_slide(prs, QR_FILE)
    prs.core_properties.title = (
        "AIは『道具』でなく『相棒』― PresenterStop 25枚版"
    )
    prs.core_properties.subject = "公開参照リンク付き・最終まとめスライド追加版"
    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    output = build()
    print(output)
    print(f"slides=25 public_references={len(REFERENCES)}")
