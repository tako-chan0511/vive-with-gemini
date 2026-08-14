#!/usr/bin/env python3
"""Repair the QR codes and public links in Agile_AI_Partnership.pptx.

The source deck is a set of full-slide PNG images.  Its visible QR codes and
footer URLs are therefore pixels rather than editable PowerPoint objects.  This
script covers the invalid pixels with deterministic replacement images and
adds real PowerPoint hyperlinks to the exact matching headings on the public
VitePress site.

Run with:

    PYTHONPATH=/tmp/agile_ppt_deps \
      python3 scripts/fix_agile_ai_partnership_links.py
"""

from __future__ import annotations

import shutil
import tempfile
import unicodedata
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from urllib.parse import quote

import qrcode
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION = ROOT / "presentation" / "Agile_AI_Partnership.pptx"

BASE = "https://tako-chan0511.github.io/vive-with-gemini/"
PAGE_5 = f"{BASE}ai-agile-vive-with-gemini-5.html"
PAGE_51 = f"{BASE}ai-agile-vive-with-gemini-5-1.html"
PAGE_52 = f"{BASE}ai-agile-vive-with-gemini-5-2.html"
PAGE_53 = f"{BASE}ai-agile-vive-with-gemini-5-3.html"

SOURCE_IMAGE_SIZE = (1376, 768)
FOOTER_BBOX = (790, 710, 1358, 740)
SLIDE1_URL_BBOX = (1098, 636, 1348, 714)
SLIDE1_QR_BBOX = (952, 598, 1102, 748)
SLIDE15_QR_BBOX = (1228, 582, 1358, 712)

FONT_FILE = Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf")


@dataclass(frozen=True)
class SlideLink:
    label: str
    page_url: str
    heading_id: str

    @property
    def target_url(self) -> str:
        fragment = unicodedata.normalize("NFD", self.heading_id)
        return f"{self.page_url}#{quote(fragment, safe='-._~')}"


LINKS: dict[int, SlideLink] = {
    1: SlideLink("AIは『道具』でなく『相棒』", PAGE_5, 'aiは「道具」でなく「相棒」'),
    2: SlideLink("自己紹介・26個のアプリ", PAGE_5, "原-桂介-keisuke-hara"),
    3: SlideLink(
        "なぜフィードバックループなのか",
        PAGE_51,
        "_1-なぜフィードバックループなのか",
    ),
    4: SlideLink("『相棒』とは何か", PAGE_51, '_2-「相棒」とは何か'),
    5: SlideLink("AI活用の3重の安全網", PAGE_52, "_1-実践の前提"),
    6: SlideLink(
        "設計から確認までのフィードバックループ",
        PAGE_52,
        "_4-aiと自動化でフィードバックを速くする",
    ),
    7: SlideLink(
        "閉域環境向けMaaS基盤",
        PAGE_53,
        "now-閉域環境向けお客様社内maas-モデル提供サービス-基盤",
    ),
    8: SlideLink(
        "AI利用領域と閉域領域の分離",
        PAGE_53,
        "_1️⃣-ai利用領域と閉域領域を分離する",
    ),
    9: SlideLink(
        "Microsoft 365 CopilotとCodexの役割",
        PAGE_53,
        "_2️⃣-microsoft-365-copilotとcodexの役割",
    ),
    10: SlideLink(
        "Kong・OpenShift／Kubernetes・GPUの責任分離",
        PAGE_53,
        "_3️⃣-kong、openshift-kubernetes、gpuの責任を分ける",
    ),
    11: SlideLink(
        "性能ベンチマーク自動化",
        PAGE_53,
        "_4️⃣-多様なモデルを短時間で比較する性能ベンチマーク自動化",
    ),
    12: SlideLink(
        "Copilot Agentによる評価ブック作成",
        PAGE_53,
        "copilot-agentによる評価ブック作成の流れ",
    ),
    13: SlideLink(
        "GitLab Issueを利用した1週間スプリント",
        PAGE_53,
        "_5️⃣-gitlab-issueを利用した1週間スプリント",
    ),
    14: SlideLink(
        "4本柱を1つのフィードバックループへ",
        PAGE_53,
        "_6️⃣-4本柱を1つのフィードバックループへ",
    ),
    15: SlideLink("結び", PAGE_53, "💬-結び"),
}


def pixel_bbox_to_emu(picture, bbox: tuple[int, int, int, int]):
    """Map a bounding box in the source PNG to the picture's PPT coordinates."""

    x0, y0, x1, y1 = bbox
    image_w, image_h = SOURCE_IMAGE_SIZE
    left = int(picture.left + picture.width * x0 / image_w)
    top = int(picture.top + picture.height * y0 / image_h)
    width = int(picture.width * (x1 - x0) / image_w)
    height = int(picture.height * (y1 - y0) / image_h)
    return left, top, width, height


def remove_generated_shapes(slide) -> None:
    """Make the repair idempotent when the script is run more than once."""

    prefixes = ("PublicLink_", "PublicURL_", "QRLink_")
    for shape in list(slide.shapes):
        if shape.name.startswith(prefixes):
            shape._element.getparent().remove(shape._element)


def add_invisible_hotspot(slide, picture, bbox, name: str, target_url: str) -> None:
    left, top, width, height = pixel_bbox_to_emu(picture, bbox)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.name = name
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # A fully transparent solid fill remains a reliable whole-box click target.
    color = shape._element.spPr.xpath("./a:solidFill/a:srgbClr")[0]
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "0")
    color.append(alpha)
    shape.click_action.hyperlink.address = target_url


def make_qr(path: Path, payload: str) -> None:
    qr = qrcode.QRCode(
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=16,
        border=4,
    )
    qr.add_data(payload)
    qr.make(fit=True)
    image = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    image.save(path, format="PNG")


def make_footer_patch(background: Image.Image, page_url: str, output: Path) -> None:
    """Remove the baked-in URL and draw one exact, readable page URL."""

    x0, y0, x1, y1 = FOOTER_BBOX
    # Derive a clean horizontal background profile from light pixels in a tall
    # band around the footer.  The tall sample and light-pixel filter exclude
    # the old URL, the Gemini mark, and nearby diagram strokes without leaving
    # a blurred "ghost" of the invalid text.
    profile: list[tuple[int, int, int]] = []
    for source_x in range(x0, x1):
        light_pixels = [
            background.getpixel((source_x, source_y))[:3]
            for source_y in range(670, SOURCE_IMAGE_SIZE[1])
            if min(background.getpixel((source_x, source_y))[:3]) >= 180
        ]
        if not light_pixels:
            light_pixels = [(242, 238, 228)]
        channels = zip(*light_pixels)
        profile.append(tuple(sorted(channel)[len(light_pixels) // 2] for channel in channels))

    strip = Image.new("RGB", (x1 - x0, 9))
    for strip_y in range(strip.height):
        for strip_x, color in enumerate(profile):
            strip.putpixel((strip_x, strip_y), color)
    strip = strip.filter(ImageFilter.GaussianBlur(radius=8))
    patch = strip.resize((x1 - x0, y1 - y0), Image.Resampling.BILINEAR)

    draw = ImageDraw.Draw(patch)
    font = ImageFont.truetype(str(FONT_FILE), 14)
    text_width = float(draw.textlength(page_url, font=font))
    text_box = draw.textbbox((0, 0), page_url, font=font)
    text_height = text_box[3] - text_box[1]
    x = max(6.0, patch.width - text_width - 8.0)
    y = (patch.height - text_height) / 2.0 - text_box[1]
    color = (19, 48, 77)
    draw.text((x, y), page_url, fill=color, font=font)
    underline_y = min(patch.height - 3, int(y + text_box[3] + 1))
    draw.line((x, underline_y, x + text_width, underline_y), fill=color, width=1)
    patch.save(output, format="PNG")


def add_linked_picture(slide, picture, source: Path, bbox, name: str, target_url: str):
    left, top, width, height = pixel_bbox_to_emu(picture, bbox)
    added = slide.shapes.add_picture(str(source), left, top, width, height)
    added.name = name
    added.click_action.hyperlink.address = target_url
    added._element.nvPicPr.cNvPr.set("descr", target_url)
    return added


def build() -> Path:
    if not PRESENTATION.exists():
        raise SystemExit(f"Presentation not found: {PRESENTATION}")
    if not FONT_FILE.exists():
        raise SystemExit(f"Footer font not found: {FONT_FILE}")

    prs = Presentation(PRESENTATION)
    if len(prs.slides) != len(LINKS):
        raise SystemExit(f"Expected {len(LINKS)} slides, found {len(prs.slides)}")

    assets = Path(tempfile.mkdtemp(prefix="agile-ai-partnership-", dir="/tmp"))
    slide1_qr = assets / "slide1-qr.png"
    slide15_qr = assets / "slide15-qr.png"
    # Keep QR payloads relatively short for robust scanning at presentation size.
    make_qr(slide1_qr, PAGE_5)
    make_qr(slide15_qr, PAGE_53)

    for slide_number, slide in enumerate(prs.slides, start=1):
        remove_generated_shapes(slide)
        if not slide.shapes or not slide.shapes[0].shape_type == 13:
            raise SystemExit(f"Slide {slide_number} does not start with a picture")
        background_picture = slide.shapes[0]
        background = Image.open(BytesIO(background_picture.image.blob)).convert("RGB")
        if background.size != SOURCE_IMAGE_SIZE:
            raise SystemExit(
                f"Slide {slide_number} image size is {background.size}, "
                f"expected {SOURCE_IMAGE_SIZE}"
            )

        link = LINKS[slide_number]
        if slide_number == 1:
            add_invisible_hotspot(
                slide,
                background_picture,
                SLIDE1_URL_BBOX,
                "PublicURL_01",
                link.target_url,
            )
            add_linked_picture(
                slide,
                background_picture,
                slide1_qr,
                SLIDE1_QR_BBOX,
                "QRLink_01",
                link.target_url,
            )
            continue

        footer = assets / f"footer-{slide_number:02d}.png"
        make_footer_patch(background, link.page_url, footer)
        add_linked_picture(
            slide,
            background_picture,
            footer,
            FOOTER_BBOX,
            f"PublicLink_{slide_number:02d}",
            link.target_url,
        )

        if slide_number == 15:
            add_linked_picture(
                slide,
                background_picture,
                slide15_qr,
                SLIDE15_QR_BBOX,
                "QRLink_15",
                link.target_url,
            )

    prs.core_properties.subject = "QR再作成・公開HP見出し別リンク修正版"

    temporary_output = PRESENTATION.with_name(
        f".{PRESENTATION.stem}.corrected.tmp{PRESENTATION.suffix}"
    )
    backup = Path("/tmp") / f"{PRESENTATION.stem}.before-link-repair.pptx"
    shutil.copy2(PRESENTATION, backup)
    prs.save(temporary_output)
    temporary_output.replace(PRESENTATION)

    print(PRESENTATION)
    print(f"slides={len(prs.slides)} precise_links={len(LINKS)} qr_codes=2")
    for number, link in LINKS.items():
        print(f"{number:02d}: {link.target_url}")
    return PRESENTATION


if __name__ == "__main__":
    build()
