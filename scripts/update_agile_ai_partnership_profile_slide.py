#!/usr/bin/env python3
"""Update the self-introduction slide in Agile_AI_Partnership.pptx.

The target slide is located by its existing ``PublicLink_02`` shape, so the
script works both before and after the three section dividers are inserted.
Only the full-slide background picture is replaced; the precise public link,
slide order, QR codes, and every other slide remain untouched.

Run after the link repair (and normally after section insertion):

    PYTHONPATH=/tmp/agile_ppt_deps \
      python3 scripts/update_agile_ai_partnership_profile_slide.py
"""

from __future__ import annotations

import hashlib
import shutil
import tempfile
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION = ROOT / "presentation" / "Agile_AI_Partnership.pptx"

IMAGE_W = 1376
IMAGE_H = 768
IMAGE_SIZE = (IMAGE_W, IMAGE_H)

REGULAR_FONT = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc")
BOLD_FONT = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc")

CREAM_TOP = (249, 247, 241)
# Match the solid fill used by the existing clickable footer-link overlay.
# Keeping the bottom edge identical avoids a visible seam after the raster
# background is replaced without touching the link shape itself.
CREAM_BOTTOM = (246, 243, 234)
NAVY = (40, 36, 87)
GOLD = (202, 148, 70)
SLATE = (81, 94, 110)
PALE_GOLD = (239, 228, 206)
NAV_LINE = (216, 201, 172)
WHITE = (255, 255, 255)


def font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(BOLD_FONT if bold else REGULAR_FONT), size)


def centered_text(draw, xy, text, text_font, fill) -> None:
    draw.text(xy, text, font=text_font, fill=fill, anchor="mm")


def base_canvas() -> Image.Image:
    image = Image.new("RGB", IMAGE_SIZE)
    draw = ImageDraw.Draw(image)
    for y in range(IMAGE_H):
        ratio = y / (IMAGE_H - 1)
        color = tuple(
            round(top + (bottom - top) * ratio)
            for top, bottom in zip(CREAM_TOP, CREAM_BOTTOM)
        )
        draw.line((0, y, IMAGE_W, y), fill=color)
    return image


def draw_keyword_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    *,
    label: str,
    title: str,
    lines: tuple[str, ...],
    active: bool = False,
) -> None:
    x0, y0, x1, y1 = box
    fill = (252, 250, 246) if not active else (250, 244, 232)
    outline = GOLD if active else NAV_LINE
    width = 4 if active else 2
    draw.rounded_rectangle(box, radius=18, fill=fill, outline=outline, width=width)
    draw.rounded_rectangle((x0 + 16, y0 + 16, x0 + 146, y0 + 46), radius=14, fill=GOLD if active else NAVY)
    centered_text(draw, (x0 + 81, y0 + 31), label, font(13, bold=True), WHITE)
    draw.text((x0 + 18, y0 + 70), title, font=font(22, bold=True), fill=NAVY)
    line_y = y0 + 116
    for line in lines:
        draw.ellipse((x0 + 19, line_y + 8, x0 + 27, line_y + 16), fill=GOLD)
        draw.text((x0 + 38, line_y), line, font=font(16, bold=active), fill=SLATE if not active else NAVY)
        line_y += 38
    if active:
        draw.rounded_rectangle((x0 + 18, y1 - 44, x1 - 18, y1 - 14), radius=14, fill=NAVY)
        centered_text(draw, ((x0 + x1) / 2, y1 - 29), "67歳・初体験に挑戦中", font(14, bold=True), WHITE)


def make_profile_image(output: Path) -> None:
    image = base_canvas()
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((64, 42, 365, 78), radius=18, fill=NAVY)
    centered_text(draw, (214, 60), "PART 01  /  SELF INTRODUCTION", font(13, bold=True), WHITE)
    draw.line((64, 98, 155, 98), fill=GOLD, width=5)

    draw.text((64, 118), "64歳からの挑戦。67歳のいまも、初挑戦。", font=font(42, bold=True), fill=NAVY)
    draw.text(
        (66, 181),
        "この3年間で、資格取得から26個のアプリへ。現在はAIモデル作成・評価とインフラ構築に挑戦中。",
        font=font(20),
        fill=SLATE,
    )

    draw.text((1010, 98), "64", font=font(72, bold=True), fill=PALE_GOLD)
    draw.text((1115, 128), "→", font=font(33, bold=True), fill=GOLD)
    draw.text((1180, 86), "67", font=font(92, bold=True), fill=NAVY)
    centered_text(draw, (1170, 195), "3年間、挑戦を継続中", font(14, bold=True), GOLD)

    card_y0 = 245
    card_y1 = 510
    cards = (
        ((64, card_y0, 352, card_y1), "LEARN", "資格取得", ("G検定 / AWS CCP", "学びを実践へ"), False),
        ((372, card_y0, 660, card_y1), "BUILD", "Vue 3 × Python", ("API Gateway", "Lambda × S3"), False),
        ((680, card_y0, 968, card_y1), "EXPAND", "K8s / OpenShift", ("クラスタ環境", "モデル性能評価"), False),
        ((988, card_y0, 1312, card_y1), "CHALLENGE NOW", "AIモデル作成・評価", ("インフラ構築", "MaaS環境構築"), True),
    )
    for box, label, title, lines, active in cards:
        draw_keyword_card(
            draw,
            box,
            label=label,
            title=title,
            lines=lines,
            active=active,
        )

    draw.rounded_rectangle((64, 544, 1312, 682), radius=20, fill=(252, 250, 246), outline=NAV_LINE, width=2)
    draw.text((88, 557), "26", font=font(66, bold=True), fill=GOLD)
    draw.text((181, 575), "APPS", font=font(25, bold=True), fill=NAVY)
    draw.line((270, 562, 270, 663), fill=NAV_LINE, width=2)
    draw.text((304, 563), "GitHub × 主要クラウド連携", font=font(24, bold=True), fill=NAVY)
    draw.rounded_rectangle((304, 613, 475, 650), radius=18, fill=GOLD)
    centered_text(draw, (390, 631), "すべて無料枠", font(15, bold=True), WHITE)
    draw.text((495, 613), "で構築・公開", font=font(18, bold=True), fill=SLATE)
    draw.line((760, 562, 760, 663), fill=NAV_LINE, width=2)
    draw.text((795, 574), "年齢ではなく、", font=font(19, bold=True), fill=SLATE)
    draw.text((795, 610), "好奇心が次の一歩を決める。", font=font(21, bold=True), fill=NAVY)
    draw.text((795, 646), "これが、私自身の Vive with Gemini。", font=font(14, bold=True), fill=GOLD)

    draw.text(
        (64, 724),
        "原 桂介 / Keisuke Hara  ｜  ビートテック株式会社 九州支店",
        font=font(12, bold=True),
        fill=SLATE,
    )
    image.save(output, format="PNG")


def find_profile_slide(prs: Presentation):
    matches = [
        slide
        for slide in prs.slides
        if any(shape.name == "PublicLink_02" for shape in slide.shapes)
    ]
    if len(matches) != 1:
        raise SystemExit(f"Expected one profile slide, found {len(matches)}")
    return matches[0]


def locate_background_media(slide) -> tuple[str, bytes]:
    pictures = [shape for shape in slide.shapes if hasattr(shape, "image")]
    backgrounds = [shape for shape in pictures if shape.name != "PublicLink_02"]
    if len(backgrounds) != 1:
        raise SystemExit(f"Expected one profile background, found {len(backgrounds)}")
    background = backgrounds[0]
    rel_id = background._element.blipFill.blip.rEmbed
    partname = str(slide.part.rels[rel_id].target_part.partname).lstrip("/")
    return partname, background.image.blob


def replace_zip_member(source: Path, output: Path, member: str, data: bytes) -> None:
    """Copy a PPTX archive while replacing exactly one uncompressed member."""
    replaced = 0
    with zipfile.ZipFile(source, "r") as src, zipfile.ZipFile(output, "w") as dst:
        for info in src.infolist():
            payload = src.read(info.filename)
            if info.filename == member:
                payload = data
                replaced += 1
            dst.writestr(info, payload)
    if replaced != 1:
        output.unlink(missing_ok=True)
        raise SystemExit(f"Expected one ZIP member named {member}, found {replaced}")


def build() -> Path:
    if not PRESENTATION.exists():
        raise SystemExit(f"Presentation not found: {PRESENTATION}")
    if not REGULAR_FONT.exists() or not BOLD_FONT.exists():
        raise SystemExit("Noto Sans CJK fonts are required")

    assets = Path(tempfile.mkdtemp(prefix="agile-ai-profile-", dir="/tmp"))
    profile_image = assets / "profile-slide.png"
    make_profile_image(profile_image)

    prs = Presentation(PRESENTATION)
    slide = find_profile_slide(prs)
    media_member, current_image = locate_background_media(slide)
    profile_bytes = profile_image.read_bytes()
    if hashlib.sha256(current_image).digest() == hashlib.sha256(profile_bytes).digest():
        print(PRESENTATION)
        print("profile_slide=already_current changed=false")
        return PRESENTATION

    backup = Path("/tmp") / f"{PRESENTATION.stem}.before-profile-update.pptx"
    temporary_output = PRESENTATION.with_name(
        f".{PRESENTATION.stem}.profile.tmp{PRESENTATION.suffix}"
    )
    shutil.copy2(PRESENTATION, backup)
    replace_zip_member(PRESENTATION, temporary_output, media_member, profile_bytes)
    temporary_output.replace(PRESENTATION)

    print(PRESENTATION)
    print(
        "profile_slide=updated age=67 experience_years=3 apps=26 "
        f"free_tier=true member={media_member}"
    )
    return PRESENTATION


if __name__ == "__main__":
    build()
