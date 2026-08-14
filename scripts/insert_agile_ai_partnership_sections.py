#!/usr/bin/env python3
"""Insert three narrative section dividers into Agile_AI_Partnership.pptx.

The repaired source deck contains 15 full-slide PNGs.  This script preserves
those slide parts verbatim, generates three matching 1376x768 divider images,
adds precise public-page and internal navigation hotspots, and reorders only
the presentation slide-id list.

Run after ``fix_agile_ai_partnership_links.py``:

    PYTHONPATH=/tmp/agile_ppt_deps \
      python3 scripts/insert_agile_ai_partnership_sections.py
"""

from __future__ import annotations

import shutil
import tempfile
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import quote

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION = ROOT / "presentation" / "Agile_AI_Partnership.pptx"

BASE = "https://tako-chan0511.github.io/vive-with-gemini/"
PAGE_5 = f"{BASE}ai-agile-vive-with-gemini-5.html"
PAGE_51 = f"{BASE}ai-agile-vive-with-gemini-5-1.html"
PAGE_53 = f"{BASE}ai-agile-vive-with-gemini-5-3.html"

IMAGE_W = 1376
IMAGE_H = 768
IMAGE_SIZE = (IMAGE_W, IMAGE_H)
PICTURE_TOP = 38100
PICTURE_WIDTH = 16256000
PICTURE_HEIGHT = 9067800

REGULAR_FONT = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc")
BOLD_FONT = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc")

CREAM_TOP = (249, 247, 241)
CREAM_BOTTOM = (242, 237, 227)
NAVY = (40, 36, 87)
GOLD = (202, 148, 70)
INK = (21, 21, 20)
SLATE = (81, 94, 110)
PALE_GOLD = (232, 220, 196)
NAV_LINE = (216, 201, 172)
WHITE = (255, 255, 255)

NAV_BBOXES = (
    (50, 590, 425, 690),
    (470, 590, 910, 690),
    (955, 590, 1330, 690),
)
PUBLIC_LINK_BBOX = (760, 704, 1358, 744)


def heading_url(page: str, heading_id: str) -> str:
    fragment = unicodedata.normalize("NFD", heading_id)
    return f"{page}#{quote(fragment, safe='-._~')}"


@dataclass(frozen=True)
class Section:
    number: int
    kicker: str
    title_lines: tuple[tuple[str, tuple[int, int, int]], ...]
    title_size: int
    description: str
    bridge: str
    page_url: str
    target_url: str
    motif: str


SECTIONS = (
    Section(
        1,
        "PART 01  /  INTRODUCTION",
        (("Vive with Geminiの思想と私", NAVY),),
        50,
        "AIを相棒として対話しながら、\n開発・学習・改善を進める実践の総称。",
        "自己紹介とともに、私が体験した2つのAI活用プロジェクトへ。",
        PAGE_5,
        heading_url(PAGE_5, 'aiは「道具」でなく「相棒」'),
        "partner",
    ),
    Section(
        2,
        "PART 02  /  PROJECT 01",
        (
            ("顧客管理システム刷新プロジェクト", NAVY),
            ("における AI駆動開発", GOLD),
        ),
        43,
        "AIを設計・実装・テストに組み込み、\n短いフィードバックループで改善を回す。",
        "Vive with Geminiの思想が、現場で自然発生アジャイルとして表れた。",
        PAGE_51,
        heading_url(PAGE_51, "プロジェクト概要"),
        "loop",
    ),
    Section(
        3,
        "PART 03  /  PROJECT 02",
        (
            ("閉域環境向け", GOLD),
            ("お客様社内MaaS基盤", NAVY),
        ),
        47,
        "Model as a Service（モデル提供サービス）\nAI・人・基盤を分離しながら、安全に短い周期で学び続ける。",
        "Vive with Geminiの同じ思想を、閉域環境の基盤構築と評価自動化へ。",
        PAGE_53,
        heading_url(
            PAGE_53,
            "now-閉域環境向けお客様社内maas-モデル提供サービス-基盤",
        ),
        "boundary",
    ),
)


def font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(BOLD_FONT if bold else REGULAR_FONT), size)


def centered_text(
    draw: ImageDraw.ImageDraw,
    xy: tuple[float, float],
    text: str,
    text_font: ImageFont.FreeTypeFont,
    fill,
) -> None:
    draw.text(xy, text, font=text_font, fill=fill, anchor="mm")


def base_canvas() -> Image.Image:
    image = Image.new("RGB", IMAGE_SIZE)
    draw = ImageDraw.Draw(image)
    for y in range(IMAGE_H):
        ratio = y / (IMAGE_H - 1)
        color = tuple(
            round(top + (bottom - top) * ratio)
            for top, bottom in zip(CREAM_TOP, CREAM_BOTTOM)
        )
        draw.line((0, y, IMAGE_W, y), fill=color)
    return image


def rotated_ellipse(
    image: Image.Image,
    box: tuple[int, int, int, int],
    angle: float,
    color,
    width: int,
) -> None:
    x0, y0, x1, y1 = box
    motif = Image.new("RGBA", (x1 - x0 + 40, y1 - y0 + 40), (0, 0, 0, 0))
    motif_draw = ImageDraw.Draw(motif)
    motif_draw.ellipse(
        (20, 20, motif.width - 20, motif.height - 20), outline=color, width=width
    )
    motif = motif.rotate(angle, resample=Image.Resampling.BICUBIC, expand=True)
    left = round((x0 + x1 - motif.width) / 2)
    top = round((y0 + y1 - motif.height) / 2)
    image.paste(motif, (left, top), motif)


def draw_partner_motif(image: Image.Image) -> None:
    rotated_ellipse(image, (1010, 348, 1270, 496), 0, NAVY, 4)
    rotated_ellipse(image, (1038, 330, 1242, 516), 55, GOLD, 4)
    rotated_ellipse(image, (1038, 330, 1242, 516), -55, NAVY, 4)


def draw_loop_motif(draw: ImageDraw.ImageDraw) -> None:
    box = (1010, 330, 1280, 530)
    draw.arc(box, 25, 190, fill=NAVY, width=6)
    draw.arc(box, 205, 365, fill=GOLD, width=6)
    draw.polygon(((1027, 355), (1050, 350), (1039, 374)), fill=NAVY)
    draw.polygon(((1262, 505), (1238, 510), (1250, 486)), fill=GOLD)
    for x, y in ((1145, 330), (1280, 430), (1145, 530), (1010, 430)):
        draw.ellipse((x - 7, y - 7, x + 7, y + 7), fill=GOLD, outline=NAVY, width=2)
    centered_text(draw, (1145, 430), "短い\nループ", font(23, bold=True), NAVY)


def draw_boundary_motif(draw: ImageDraw.ImageDraw) -> None:
    draw.rounded_rectangle((970, 365, 1090, 475), radius=14, outline=GOLD, width=4)
    centered_text(draw, (1030, 420), "Cloud\nAI", font(20, bold=True), GOLD)
    draw.line((1135, 335, 1135, 515), fill=NAVY, width=6)
    draw.rounded_rectangle((1114, 413, 1156, 455), radius=8, fill=CREAM_BOTTOM, outline=NAVY, width=3)
    draw.arc((1123, 392, 1147, 428), 180, 360, fill=NAVY, width=3)
    draw.rounded_rectangle((1180, 365, 1320, 475), radius=14, outline=NAVY, width=4)
    centered_text(draw, (1250, 420), "社内\nMaaS", font(20, bold=True), NAVY)
    centered_text(draw, (1145, 498), "人のレビュー・承認", font(14, bold=True), SLATE)


def draw_navigation(draw: ImageDraw.ImageDraw, active: int) -> None:
    nodes = (198, 688, 1178)
    y = 626
    draw.line((nodes[0], y, nodes[-1], y), fill=NAV_LINE, width=3)
    if active >= 2:
        draw.line((nodes[0], y, nodes[1], y), fill=NAVY, width=4)
    if active >= 3:
        draw.line((nodes[1], y, nodes[2], y), fill=NAVY, width=4)

    labels = (
        "01  思想・自己紹介",
        "02  顧客管理 × AI駆動開発",
        "03  閉域MaaS基盤",
    )
    label_x = (235, 690, 1140)
    for index, (node_x, label, center_x) in enumerate(zip(nodes, labels, label_x), 1):
        if index < active:
            radius = 10
            draw.ellipse(
                (node_x - radius, y - radius, node_x + radius, y + radius),
                fill=NAVY,
            )
            label_color = NAVY
            label_font = font(16)
        elif index == active:
            radius = 16
            draw.ellipse(
                (node_x - radius, y - radius, node_x + radius, y + radius),
                fill=GOLD,
                outline=NAVY,
                width=2,
            )
            label_color = NAVY
            label_font = font(17, bold=True)
        else:
            radius = 10
            draw.ellipse(
                (node_x - radius, y - radius, node_x + radius, y + radius),
                fill=CREAM_BOTTOM,
                outline=(154, 149, 138),
                width=2,
            )
            label_color = (125, 121, 112)
            label_font = font(16)
        centered_text(draw, (center_x, 665), label, label_font, label_color)


def make_section_image(section: Section, output: Path) -> None:
    image = base_canvas()
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((64, 43, 350, 78), radius=17, fill=NAVY)
    centered_text(draw, (207, 60), section.kicker, font(14, bold=True), WHITE)
    draw.line((64, 98, 155, 98), fill=GOLD, width=5)

    draw.text((1000, 68), f"{section.number:02d}", font=font(190, bold=True), fill=PALE_GOLD)
    centered_text(draw, (1290, 47), f"{section.number:02d} / 03", font(14, bold=True), SLATE)

    title_y = 128
    title_font = font(section.title_size, bold=True)
    line_step = section.title_size + 18
    for line_index, (line, color) in enumerate(section.title_lines):
        draw.text((64, title_y + line_index * line_step), line, font=title_font, fill=color)

    draw.multiline_text(
        (66, 320),
        section.description,
        font=font(25),
        fill=SLATE,
        spacing=13,
    )
    draw.text((66, 438), section.bridge, font=font(18, bold=True), fill=GOLD)

    draw.rounded_rectangle((64, 490, 920, 552), radius=14, fill=(252, 250, 246), outline=NAV_LINE, width=2)
    draw.rounded_rectangle((80, 506, 255, 536), radius=14, fill=GOLD)
    centered_text(draw, (168, 521), "VIVE 共通ループ", font(14, bold=True), WHITE)
    draw.text(
        (278, 503),
        "AIが提案  →  人が判断  →  小さく検証  →  次の改善",
        font=font(17, bold=True),
        fill=NAVY,
    )

    if section.motif == "partner":
        draw_partner_motif(image)
        draw = ImageDraw.Draw(image)
    elif section.motif == "loop":
        draw_loop_motif(draw)
    else:
        draw_boundary_motif(draw)

    draw_navigation(draw, section.number)
    draw.text((64, 724), "Vive with Gemini", font=font(12, bold=True), fill=SLATE)
    visible_url = f"公開HP  ↗  {section.page_url}"
    url_font = font(11)
    url_width = draw.textlength(visible_url, font=url_font)
    url_x = 1348 - url_width
    draw.text((url_x, 718), visible_url, font=url_font, fill=NAVY)
    draw.line((url_x, 735, 1348, 735), fill=NAVY, width=1)
    image.save(output, format="PNG")


def is_section_slide(slide) -> bool:
    return any(shape.name.startswith("SectionDivider_") for shape in slide.shapes)


def pixel_bbox_to_emu(picture, bbox: tuple[int, int, int, int]):
    x0, y0, x1, y1 = bbox
    left = int(picture.left + picture.width * x0 / IMAGE_W)
    top = int(picture.top + picture.height * y0 / IMAGE_H)
    width = int(picture.width * (x1 - x0) / IMAGE_W)
    height = int(picture.height * (y1 - y0) / IMAGE_H)
    return left, top, width, height


def add_hotspot(slide, picture, bbox, name: str):
    left, top, width, height = pixel_bbox_to_emu(picture, bbox)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.name = name
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    color = shape._element.spPr.xpath("./a:solidFill/a:srgbClr")[0]
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "0")
    color.append(alpha)
    return shape


def add_section_slide(prs: Presentation, section: Section, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*CREAM_BOTTOM)
    picture = slide.shapes.add_picture(
        str(image_path), 0, PICTURE_TOP, PICTURE_WIDTH, PICTURE_HEIGHT
    )
    picture.name = f"SectionDivider_{section.number:02d}_Background"
    picture._element.nvPicPr.cNvPr.set(
        "descr", f"{section.kicker}: {' / '.join(line for line, _ in section.title_lines)}"
    )
    public_link = add_hotspot(
        slide,
        picture,
        PUBLIC_LINK_BBOX,
        f"SectionDivider_{section.number:02d}_PublicLink",
    )
    public_link.click_action.hyperlink.address = section.target_url
    nav_hotspots = [
        add_hotspot(
            slide,
            picture,
            bbox,
            f"SectionDivider_{section.number:02d}_Nav_{nav_number:02d}",
        )
        for nav_number, bbox in enumerate(NAV_BBOXES, 1)
    ]
    return slide, nav_hotspots


def reorder_slides(prs: Presentation, original_ids, section_ids) -> None:
    # Cover -> Part 1 -> self intro -> Part 2 -> CRM slides -> Part 3 -> MaaS -> close
    final_order = [
        original_ids[0],
        section_ids[0],
        original_ids[1],
        section_ids[1],
        *original_ids[2:6],
        section_ids[2],
        *original_ids[6:],
    ]
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_id in final_order:
        slide_id_list.append(slide_id)


def build() -> Path:
    if not PRESENTATION.exists():
        raise SystemExit(f"Presentation not found: {PRESENTATION}")
    if not REGULAR_FONT.exists() or not BOLD_FONT.exists():
        raise SystemExit("Noto Sans CJK fonts are required to render section dividers")

    prs = Presentation(PRESENTATION)
    existing_sections = [slide for slide in prs.slides if is_section_slide(slide)]
    if len(prs.slides) == 18 and len(existing_sections) == 3:
        print(PRESENTATION)
        print("slides=18 section_dividers=3 already_present=true")
        return PRESENTATION
    if existing_sections:
        raise SystemExit(
            f"Found {len(existing_sections)} section dividers in an unexpected deck"
        )
    if len(prs.slides) != 15:
        raise SystemExit(f"Expected 15 content slides, found {len(prs.slides)}")
    if prs.slide_width != PICTURE_WIDTH or prs.slide_height != 9144000:
        raise SystemExit("Unexpected presentation dimensions")

    original_ids = list(prs.slides._sldIdLst)
    assets = Path(tempfile.mkdtemp(prefix="agile-ai-sections-", dir="/tmp"))
    section_slides = []
    section_ids = []
    section_navs = []
    for section in SECTIONS:
        image_path = assets / f"section-{section.number:02d}.png"
        make_section_image(section, image_path)
        slide, navs = add_section_slide(prs, section, image_path)
        section_slides.append(slide)
        section_navs.append(navs)
        section_ids.append(list(prs.slides._sldIdLst)[-1])

    for navs in section_navs:
        for hotspot, target_slide in zip(navs, section_slides):
            hotspot.click_action.target_slide = target_slide

    reorder_slides(prs, original_ids, section_ids)
    prs.core_properties.subject = "QR・精密リンク修正版／3部構成セクション扉付き"

    backup = Path("/tmp") / f"{PRESENTATION.stem}.before-section-dividers.pptx"
    temporary_output = PRESENTATION.with_name(
        f".{PRESENTATION.stem}.sections.tmp{PRESENTATION.suffix}"
    )
    shutil.copy2(PRESENTATION, backup)
    prs.save(temporary_output)
    temporary_output.replace(PRESENTATION)

    print(PRESENTATION)
    print("slides=18 section_dividers=3 internal_nav_links=9 public_links=3")
    for section in SECTIONS:
        print(f"section{section.number:02d}: {section.target_url}")
    return PRESENTATION


if __name__ == "__main__":
    build()
